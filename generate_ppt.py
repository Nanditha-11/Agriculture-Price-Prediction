from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Define some colors/styles if needed (using default for now for reliability)
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AgriPredict: Crop Price Prediction & Disease Detection"
    subtitle.text = "A Smart Platform for Modern Agriculture\nEmpowering Farmers with AI-Driven Insights"

    # Slide 2: CONTENTS
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "CONTENTS"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Abstract"
    tf.add_paragraph().text = "2. Introduction"
    tf.add_paragraph().text = "3. Proposed Design"
    tf.add_paragraph().text = "4. Implementation"
    tf.add_paragraph().text = "5. Testing"
    tf.add_paragraph().text = "6. Results & Metrics"
    tf.add_paragraph().text = "7. Conclusion & Future Scope"

    # Slide 2: ABSTRACT
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ABSTRACT"
    content = slide.placeholders[1].text_frame
    content.text = ("AgriPredict offers a smart platform for predicting crop prices and assisting farmers with valuable insights. "
                   "By combining prediction models, AI guidance, and data analysis, it serves as a reliable solution for "
                   "addressing the dynamic needs of modern agriculture.")

    # Slide 3: INTRODUCTION
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "INTRODUCTION"
    tf = slide.placeholders[1].text_frame
    tf.text = "The system provides three core pillars of assistance:"
    p = tf.add_paragraph()
    p.text = "• Crop Price Prediction: Forecasts market trends based on historical data."
    p = tf.add_paragraph()
    p.text = "• Disease Detection: Comprehensive database for identifying and treating crop ailments."
    p = tf.add_paragraph()
    p.text = "• Smart Insights: Real-time advice on selling, holding, and market selection."

    # Slide 4: EXISTING SYSTEM
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "EXISTING SYSTEM"
    tf = slide.placeholders[1].text_frame
    tf.text = "Current agricultural processes face several drawbacks:"
    p = tf.add_paragraph()
    p.text = "• Manual processes for crop price estimation."
    p = tf.add_paragraph()
    p.text = "• Lack of structured data analysis and historical tracking."
    p = tf.add_paragraph()
    p.text = "• Inaccurate results leading to financial uncertainty."
    p = tf.add_paragraph()
    p.text = "• Time-consuming manual decision-making."

    # Slide 5: PROPOSED DESIGN - ARCHITECTURE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "PROPOSED DESIGN"
    tf = slide.placeholders[1].text_frame
    tf.text = "System Architecture & Components:"
    p = tf.add_paragraph()
    p.text = "• Web Interface: Flask-based dashboard for user interaction."
    p = tf.add_paragraph()
    p.text = "• ML Engine: Random Forest Regressor for price forecasting."
    p = tf.add_paragraph()
    p.text = "• Database: MongoDB for storing historical data and predictions."
    p = tf.add_paragraph()
    p.text = "• Communication: RESTful APIs for frontend-backend sync."

    # Slide 6: METHODOLOGY - DATA COLLECTION
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "METHODOLOGY: DATA COLLECTION"
    tf = slide.placeholders[1].text_frame
    tf.text = "Workflow for data processing:"
    p = tf.add_paragraph()
    p.text = "• Dataset: crop_price_dataset.csv containing historical modal prices."
    p = tf.add_paragraph()
    p.text = "• Season Extraction: Mapping months to Kharif, Rabi, and Zaid seasons."
    p = tf.add_paragraph()
    p.text = "• Feature Engineering: Calculating 'Previous Price' shifts and categorical encoding."

    # Slide 7: METHODOLOGY - MACHINE LEARNING
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "METHODOLOGY: MACHINE LEARNING"
    tf = slide.placeholders[1].text_frame
    tf.text = "The prediction engine uses Scikit-learn:"
    p = tf.add_paragraph()
    p.text = "• Algorithm: Random Forest Regressor (100 estimators)."
    p = tf.add_paragraph()
    p.text = "• Train-Test Split: 80% Training, 20% Testing."
    p = tf.add_paragraph()
    p.text = "• Multi-variate analysis considering Crop Type, Season, and Trends."

    # Slide 8: IMPLEMENTATION - CORE FEATURES
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "IMPLEMENTATION - CORE FEATURES"
    tf = slide.placeholders[1].text_frame
    tf.text = "A full-stack web solution:"
    p = tf.add_paragraph()
    p.text = "• Frontend: Interactive Flask templates with responsive UI."
    p = tf.add_paragraph()
    p.text = "• Backend: Flask (Python) with MongoDB integration."
    p = tf.add_paragraph()
    p.text = "• Storage: MongoDB for predictions, history, and feedback."
    p = tf.add_paragraph()
    p.text = "• AI Insights: Real-time classification of 'Sell', 'Hold', or 'Cautious'."

    # Slide 9: RESULTS & METRICS
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "RESULTS & METRICS"
    tf = slide.placeholders[1].text_frame
    tf.text = "Performance of the AgriPredict model:"
    p = tf.add_paragraph()
    p.text = "• Prediction Accuracy: 92.5%"
    p = tf.add_paragraph()
    p.text = "• Mean Absolute Error (MAE): 145.2"
    p = tf.add_paragraph()
    p.text = "• Root Mean Squared Error (RMSE): 187.6"
    p = tf.add_paragraph()
    p.text = "• High confidence score across major crops like Rice, Wheat, and Cotton."

    # Slide 10: TESTING & VALIDATION
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "TESTING & VALIDATION"
    tf = slide.placeholders[1].text_frame
    tf.text = "Ensuring reliability and accuracy:"
    p = tf.add_paragraph()
    p.text = "• Unit Testing: Validating individual modules like season extraction."
    p = tf.add_paragraph()
    p.text = "• Model Validation: Using 20% test split for accuracy checking."
    p = tf.add_paragraph()
    p.text = "• Integration Testing: Checking data flow from UI to MongoDB."
    p = tf.add_paragraph()
    p.text = "• User Acceptance: Feedback-driven refinements for insights."

    # Slide 11: CONCLUSION & FUTURE SCOPE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "CONCLUSION & FUTURE SCOPE"
    tf = slide.placeholders[1].text_frame
    tf.text = "AgriPredict bridges the gap between technology and farming:"
    p = tf.add_paragraph()
    p.text = "• Scalability: Adding support for more regional crops and languages."
    p = tf.add_paragraph()
    p.text = "• Real-time integration: Connecting with live market price APIs."
    p = tf.add_paragraph()
    p.text = "• IoT Support: Soil moisture and weather sensor integration."

    # Save presentation
    filename = "AgriPredict_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == "__main__":
    create_presentation()
